from __future__ import annotations
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from models import Slide

# Layout indices in blank pptx (0=Title Slide, 1=Title and Content, 5=Blank)
_LAYOUT_MAP = {
    "title-only": 0,
    "title-bullets": 1,
    "title-image": 1,
    "blank": 5,
}

_THEME_BG = {
    "blue": RGBColor(0x1E, 0x40, 0xAF),
    "dark": RGBColor(0x1F, 0x29, 0x37),
    "minimal": RGBColor(0xFF, 0xFF, 0xFF),
}

_THEME_FG = {
    "blue": RGBColor(0xFF, 0xFF, 0xFF),
    "dark": RGBColor(0xFF, 0xFF, 0xFF),
    "minimal": RGBColor(0x11, 0x18, 0x27),
}


def build_pptx(slides: list[Slide], title: str = "Presentation") -> bytes:
    prs = Presentation()
    prs.core_properties.title = title

    for slide in slides:
        layout_idx = _LAYOUT_MAP.get(slide.layout, 1)
        # Clamp to available layouts
        layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
        sl = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set background colour
        bg_color = _THEME_BG.get(slide.style.theme, RGBColor(0xFF, 0xFF, 0xFF))
        fill = sl.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        fg_color = _THEME_FG.get(slide.style.theme, RGBColor(0x00, 0x00, 0x00))

        # Title placeholder
        if sl.shapes.title:
            sl.shapes.title.text = slide.title
            try:
                sl.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = fg_color
                sl.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
            except IndexError:
                pass  # no runs if title is empty

        # Body / bullets
        if slide.bullets and slide.layout in ("title-bullets", "title-image"):
            body_ph = None
            for shape in sl.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_ph = shape
                    break
            if body_ph:
                tf = body_ph.text_frame
                tf.clear()
                for i, bullet in enumerate(slide.bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    try:
                        p.runs[0].font.color.rgb = fg_color
                        p.runs[0].font.size = Pt(20)
                    except IndexError:
                        pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
