import io
import pytest
from pptx import Presentation
from services.pptx_builder import build_pptx
from models import Slide


def _make_slide(layout="title-bullets", bullets=None):
    return Slide(
        id="s1",
        title="Test Slide",
        bullets=bullets or ["Bullet 1", "Bullet 2"],
        layout=layout,
        style={"theme": "blue", "font": "sans"},
        prompt="test prompt"
    )


def test_build_pptx_returns_bytes():
    slides = [_make_slide()]
    result = build_pptx(slides, title="Test")
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_build_pptx_valid_pptx():
    slides = [_make_slide()]
    result = build_pptx(slides, title="Test")
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 1


def test_build_pptx_multiple_slides():
    slides = [_make_slide(), _make_slide(layout="title-only", bullets=[])]
    result = build_pptx(slides, title="Multi")
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 2


def test_build_pptx_title_only_layout():
    slide = _make_slide(layout="title-only", bullets=[])
    result = build_pptx([slide], title="TitleOnly")
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 1


def test_build_pptx_blank_layout():
    slide = _make_slide(layout="blank", bullets=[])
    result = build_pptx([slide], title="Blank")
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 1
